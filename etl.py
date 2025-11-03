#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
ETL - Extractor de Texto de Documentos
Soporta:
- Google Drive URLs (Docs, Sheets, Slides, archivos)
- Firebase Storage URLs
- Archivos locales

Uso:
  python etl.py --url "URL_DE_GOOGLE_DRIVE"
  python etl.py --firebase-url "gs://bucket/path/file.pdf"
  python etl.py --local-file "ruta/archivo.pdf"
"""

import os
import re
import sys
import tempfile
import argparse
import logging
import requests
import chardet
import pdfplumber
import docx
import pandas as pd
from pptx import Presentation
from urllib.parse import urlparse

# ---------- Configuración de logs ----------
def setup_logging(verbose=True):
    """Configura los logs de consola en modo DEBUG (detallado) o INFO (simple)."""
    lvl = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=lvl, format="%(asctime)s [%(levelname)s] %(message)s", datefmt="%H:%M:%S"
    )

# ---------- Identificación de URL de Google Drive ----------
def parse_drive_url(url: str):
    """
    Detecta qué tipo de recurso de Google es la URL.
    Puede ser: Docs, Sheets, Slides o archivo subido a Drive.
    Retorna una tupla (tipo, id).
    """
    # Detectar Docs, Sheets, Slides
    for t, pat in {
        "document": r"https?://docs\.google\.com/document/d/([A-Za-z0-9_\-]+)/",
        "spreadsheets": r"https?://docs\.google\.com/spreadsheets/d/([A-Za-z0-9_\-]+)/",
        "presentation": r"https?://docs\.google\.com/presentation/d/([A-Za-z0-9_\-]+)/",
    }.items():
        m = re.search(pat, url)
        if m:
            return t, m.group(1)

    # Detectar archivos subidos a Drive
    for pat in [
        r"https?://drive\.google\.com/file/d/([A-Za-z0-9_\-]+)/",
        r"https?://drive\.google\.com/open\?id=([A-Za-z0-9_\-]+)",
        r"https?://drive\.google\.com/uc\?.*?id=([A-Za-z0-9_\-]+)",
    ]:
        m = re.search(pat, url)
        if m:
            return "file", m.group(1)

    raise ValueError("URL de Google Drive/Docs no reconocida.")

def docs_export_url(doc_type, file_id):
    """Convierte Docs/Sheets/Slides a una URL de exportación directa."""
    if doc_type == "document":
        return f"https://docs.google.com/document/d/{file_id}/export?format=txt", ".txt"
    if doc_type == "spreadsheets":
        return f"https://docs.google.com/spreadsheets/d/{file_id}/export?format=xlsx", ".xlsx"
    if doc_type == "presentation":
        return f"https://docs.google.com/presentation/d/{file_id}/export?format=pptx", ".pptx"
    raise ValueError("Tipo no soportado.")

# ---------- Descarga ----------
def download_file(url, dst_path):
    """Descarga un archivo desde la URL dada y lo guarda en `dst_path`."""
    r = requests.get(url, stream=True, timeout=60)
    r.raise_for_status()
    with open(dst_path, "wb") as f:
        for chunk in r.iter_content(1024 * 1024):
            f.write(chunk)
    return dst_path

# ---------- Normalización ----------
def normalize_text(s: str):
    """Limpia texto: quita caracteres invisibles y normaliza saltos de línea."""
    s = s.replace("\u200b", "")
    lines = [re.sub(r"[ \t]+", " ", ln).rstrip() for ln in s.splitlines()]
    joined = "\n".join(lines)
    return re.sub(r"\n{3,}", "\n\n", joined).strip() + "\n"

# ---------- Extractores ----------
def extract_pdf(path):
    """Extrae texto de un PDF (si es texto digital)."""
    text = []
    with pdfplumber.open(path) as pdf:
        for i, p in enumerate(pdf.pages):
            try:
                t = p.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                text.append(f"\n=== Página {i+1} ===\n{t.strip()}")
            else:
                text.append(f"\n=== Página {i+1} ===\n(No se detectó texto en esta página)")
    return normalize_text("\n".join(text))

def extract_docx(path):
    """Extrae texto y tablas de un documento Word (DOCX)."""
    d = docx.Document(path)
    parts = [p.text.strip() for p in d.paragraphs if p.text.strip()]
    for ti, tbl in enumerate(d.tables, 1):
        parts.append(f"\n--- Tabla {ti} ---")
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
    return normalize_text("\n".join(parts))

def extract_xlsx(path):
    """Extrae texto de todas las hojas de un Excel (XLSX)."""
    out = []
    sheets = pd.read_excel(path, sheet_name=None, dtype=str, engine="openpyxl")
    for i, (name, df) in enumerate(sheets.items(), 1):
        out.append(f"\n=== Hoja {i}: {name} ===\n")
        df = df.fillna("")
        out.append(df.to_csv(index=False, sep="|"))
    return normalize_text("\n".join(out))

def extract_pptx(path):
    """Extrae títulos y textos de un PowerPoint (PPTX)."""
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n=== Slide {i} ===")
        if slide.shapes.title and slide.shapes.title.text:
            out.append(f"# {slide.shapes.title.text.strip()}")
        for shp in slide.shapes:
            if hasattr(shp, "text") and shp.text:
                t = shp.text.strip()
                if t:
                    out.append(t)
    return normalize_text("\n".join(out))

def extract_txt(path):
    """Lee un archivo de texto plano, detectando codificación automáticamente."""
    with open(path, "rb") as f:
        raw = f.read()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    return normalize_text(raw.decode(enc, errors="replace"))

def extract_image(path):
    """Ignora imágenes y devuelve un aviso en texto."""
    return "(Imagen detectada, OCR deshabilitado)\n"

def choose_extractor(path):
    """Decide qué extractor usar según la extensión del archivo."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext in [".xlsx", ".xlsm"]:
        return extract_xlsx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext in [".csv", ".tsv"]:
        return extract_xlsx(path)
    if ext in [".txt", ".md", ".log"]:
        return extract_txt(path)
    if ext in [".png", ".jpg", ".jpeg"]:
        return extract_image(path)
    return extract_txt(path)

# ---------- Firebase Support ----------
def download_from_firebase(gs_url, dst_path):
    """
    Descarga un archivo desde Firebase Storage usando la URL gs://
    Requiere que firebase_admin esté configurado.
    """
    try:
        from firebase_admin import storage
        # Extraer bucket y path de gs://bucket/path/file.pdf
        parsed = urlparse(gs_url)
        bucket_name = parsed.netloc
        file_path = parsed.path.lstrip('/')
        
        bucket = storage.bucket(bucket_name)
        blob = bucket.blob(file_path)
        blob.download_to_filename(dst_path)
        logging.info(f"Descargado desde Firebase: {file_path}")
        return dst_path
    except ImportError:
        raise ImportError("firebase_admin no está instalado. Usa: pip install firebase-admin")
    except Exception as e:
        raise Exception(f"Error descargando desde Firebase: {e}")

def download_from_http_url(url, dst_path):
    """Descarga archivo desde URL HTTP/HTTPS directa (ej: Firebase Storage public URL)"""
    try:
        r = requests.get(url, stream=True, timeout=60)
        r.raise_for_status()
        with open(dst_path, "wb") as f:
            for chunk in r.iter_content(1024 * 1024):
                f.write(chunk)
        logging.info(f"Descargado desde URL: {url}")
        return dst_path
    except Exception as e:
        raise Exception(f"Error descargando desde URL: {e}")

# ---------- Pipeline principal ----------
def process_url(url, out_path=None, verbose=True):
    """
    Flujo principal:
    1. Detecta tipo de archivo.
    2. Descarga en carpeta temporal.
    3. Extrae texto según formato.
    4. Guarda salida en .txt.
    """
    setup_logging(verbose)
    
    with tempfile.TemporaryDirectory() as td:
        # Detectar si es Firebase Storage URL
        if url.startswith("gs://"):
            logging.info("Detectado: Firebase Storage URL")
            filename = os.path.basename(urlparse(url).path)
            local = os.path.join(td, filename)
            download_from_firebase(url, local)
        
        # Detectar si es URL HTTP directa (Firebase Storage public URL)
        elif url.startswith("http") and "firebasestorage" in url:
            logging.info("Detectado: Firebase Storage HTTP URL")
            filename = "downloaded_file"
            # Intentar obtener extensión de la URL
            parsed_url = urlparse(url)
            path_parts = parsed_url.path.split('/')
            if path_parts:
                filename = path_parts[-1].split('?')[0]
            local = os.path.join(td, filename)
            download_from_http_url(url, local)
        
        # URLs de Google Drive
        else:
            kind, fid = parse_drive_url(url)
            logging.info(f"Detectado tipo={kind} id={fid}")
            
            # Si es Docs/Sheets/Slides → exportar
            if kind in ("document", "spreadsheets", "presentation"):
                u, ext = docs_export_url(kind, fid)
                local = os.path.join(td, f"{fid}{ext}")
                download_file(u, local)
            else:
                # Archivos subidos a Drive
                u = f"https://drive.google.com/uc?export=download&id={fid}"
                local = os.path.join(td, fid)
                download_file(u, local)

        # Nombre del archivo de salida
        if not out_path:
            out_path = os.path.splitext(os.path.basename(local))[0] + ".txt"

        logging.info("Extrayendo texto…")
        text = choose_extractor(local)

        with open(out_path, "w", encoding="utf-8") as f:
            f.write(text)

        logging.info(f"Listo → {out_path}")
        return out_path

def process_local_file(file_path, out_path=None, verbose=True):
    """
    Procesa un archivo local directamente sin descarga.
    Útil para archivos ya descargados de Firebase.
    """
    setup_logging(verbose)
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Archivo no encontrado: {file_path}")
    
    logging.info(f"Procesando archivo local: {file_path}")
    
    # Nombre del archivo de salida
    if not out_path:
        out_path = os.path.splitext(file_path)[0] + "_procesado.txt"
    
    logging.info("Extrayendo texto…")
    text = choose_extractor(file_path)
    
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    
    logging.info(f"Listo → {out_path}")
    return out_path, text

# ---------- CLI ----------
if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="ETL - Extractor de Texto de Documentos")
    
    # Grupo de argumentos mutuamente excluyentes
    input_group = ap.add_mutually_exclusive_group(required=True)
    input_group.add_argument("--url", help="URL pública de Google Drive")
    input_group.add_argument("--firebase-url", help="URL de Firebase Storage (gs://...)")
    input_group.add_argument("--local-file", help="Ruta a archivo local")
    
    ap.add_argument("--out", help="Ruta de salida .txt")
    ap.add_argument("--quiet", action="store_true", help="Menos logs")
    args = ap.parse_args()

    try:
        if args.local_file:
            process_local_file(args.local_file, args.out, not args.quiet)
        else:
            # Para URLs de Google Drive o Firebase
            url = args.url or args.firebase_url
            process_url(url, args.out, not args.quiet)
    except Exception as e:
        logging.exception(f"Error: {e}")
        sys.exit(1)
